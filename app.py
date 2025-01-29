import streamlit as st
import google.generativeai as genai
import os
import PyPDF2 as pdf
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
import json
import io
import time

# Load environment variables
load_dotenv()

# Configure Gemini API
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# Function to get response from Gemini API
def get_gemini_response(input_text):
    model = genai.GenerativeModel('gemini-pro')
    try:
        if not input_text.strip():
            st.error("Input text is empty. Cannot generate response.")
            return "{}"  # Return empty JSON

        response = model.generate_content(input_text)
        if response and response.text:
            return response.text
        else:
            st.error("Received an empty response from the model.")
            return "{}"  # Return empty JSON
    except Exception as e:
        st.error(f"Error while getting response from API: {str(e)}")
        return "{}"  # Return empty JSON

# Function to extract text from uploaded PDF file
def input_pdf_text(uploaded_file):
    reader = pdf.PdfReader(uploaded_file)
    text = []
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text.append(page.extract_text() or "")
    return text

# Function to extract text from uploaded Word document
def input_word_text(uploaded_file):
    doc = Document(uploaded_file)
    text = []
    for para in doc.paragraphs:
        text.append(para.text or "")
    return text

# Function to extract text from uploaded PPT file
def input_ppt_text(uploaded_file):
    presentation = Presentation(uploaded_file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text or "")
    return text

# Function to generate DOCX file from content
def generate_docx(generated_content):
    doc = Document()
    doc.add_heading('Pidgin Slide Creator Output', level=1)

    for chapter, content in generated_content.items():
        doc.add_heading(f'Chapter: {chapter}', level=2)
        for item in content:
            doc.add_paragraph(f"**Pidgin Translation:**\n{item.get('Pidgin Translation', 'No translation available.')}")
            doc.add_paragraph(f"**Questions and Answers:**\n{item.get('Questions and Answers', 'No questions available.')}")
            doc.add_paragraph()  # Add a blank line for spacing

    byte_io = io.BytesIO()
    doc.save(byte_io)
    byte_io.seek(0)
    return byte_io

# Function to group content into chapters
def group_into_chapters(text):
    chapters = {}
    current_chapter = "Chapter 1"  # Default chapter name
    chapters[current_chapter] = []

    for paragraph in text:
        if "chapter" in paragraph.lower() or "section" in paragraph.lower():
            current_chapter = paragraph.strip()
            chapters[current_chapter] = []
        else:
            chapters[current_chapter].append(paragraph)
    return chapters

# Streamlit App
st.set_page_config(page_title="Pidgin Slide Creator")
st.title("Pidgin Slide Creator")

# File uploader for slides (PDF, Word, PPT, or text) input
uploaded_file = st.file_uploader("Upload Your Document (PDF, DOCX, PPTX, TXT)...", type=["pdf", "docx", "pptx", "txt"])

# Initialize session state for history
if 'history' not in st.session_state:
    st.session_state.history = {}

# Submit button for processing the document
submit = st.button("Submit")

if uploaded_file:
    # Extract text from the uploaded file
    if uploaded_file.type == "application/pdf":
        document_text = input_pdf_text(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        document_text = input_word_text(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        document_text = input_ppt_text(uploaded_file)
    elif uploaded_file.type == "text/plain":
        document_text = uploaded_file.read().decode("utf-8").split('\n')
    else:
        st.error("Unsupported file type!")
        st.stop()

    # Group content into chapters
    chapters = group_into_chapters(document_text)

    # Allow user to select chapters
    selected_chapters = st.multiselect(
        "Select Chapters to Process",
        options=list(chapters.keys()),
        default=list(chapters.keys())  # Select all chapters by default
    )

    if submit:
        try:
            # Process selected chapters
            generated_content = {}

            # Progress bar
            progress_bar = st.progress(0)
            total_paragraphs = sum(len(chapters[chapter]) for chapter in selected_chapters)
            processed_paragraphs = 0

            for chapter in selected_chapters:
                st.markdown(f"### {chapter}")
                generated_content[chapter] = []

                for paragraph in chapters[chapter]:
                    if paragraph.strip():  # Skip empty paragraphs
                        # Prepare prompt with extracted paragraph text
                        input_prompt = """
                        You are an expert in translating and explaining content in Nigerian Pidgin English. Your task is to:
                        1. Translate the given paragraph into Nigerian Pidgin English.
                        2. Identify and answer any questions, examples, or tests in the text.
                        3. Provide detailed explanations for the answers in Nigerian Pidgin English.
                        
                        Paragraph: {paragraph}
                        
                        I want the response in the following structured format:
                        {{
                            "Pidgin Translation": "",
                            "Questions and Answers": ""
                        }}
                        """
                        input_prompt_filled = input_prompt.format(paragraph=paragraph)

                        # Get response from Gemini API
                        response = get_gemini_response(input_prompt_filled)

                        try:
                            # Parse response
                            response_json = json.loads(response)

                            # Display and collect Pidgin Translation, Questions and Answers
                            pidgin_translation = response_json.get("Pidgin Translation", "No translation available.")
                            questions_and_answers = response_json.get("Questions and Answers", "No questions available.")

                            st.markdown("**Pidgin Translation:**")
                            st.write(pidgin_translation)

                            st.markdown("**Questions and Answers:**")
                            st.write(questions_and_answers)

                            # Collect generated content in JSON format
                            paragraph_content_json = {
                                "Pidgin Translation": pidgin_translation,
                                "Questions and Answers": questions_and_answers,
                            }
                            generated_content[chapter].append(paragraph_content_json)

                            # Update progress
                            processed_paragraphs += 1
                            progress_bar.progress(processed_paragraphs / total_paragraphs)

                            # Add a 10-second delay to avoid overloading the API
                            time.sleep(10)
                        except json.JSONDecodeError:
                            st.error("Failed to decode JSON response from the model.")

            # Add generated content to history
            st.session_state.history = generated_content

            # Provide download option for DOCX
            if generated_content:
                docx_file = generate_docx(generated_content)
                st.download_button(
                    label="Download Generated Content as DOCX",
                    data=docx_file,
                    file_name='pidgin_slide_creator_output.docx',
                    mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                )
        except Exception as e:
            st.error(f"An error occurred while processing the file: {str(e)}")

# History Page
if st.sidebar.button("View History"):
    st.title("History")

    # Display history
    if 'history' in st.session_state:
        for chapter, content in st.session_state.history.items():
            st.subheader(f"{chapter}")
            for item in content:
                st.markdown("**Pidgin Translation:**")
                st.write(item.get("Pidgin Translation", "No translation available."))
                st.markdown("**Questions and Answers:**")
                st.write(item.get("Questions and Answers", "No questions available."))
                st.write("---")
    else:
        st.write("No history available.")

st.markdown("""
<style>
[data-testid="stAppViewContainer"] {
    background-image: url("https://images.unsplash.com/photo-1698945746290-a9d1cc575e77");
    background-size: cover;
    background-repeat: no-repeat;
    background-position: center;
}
</style>
""", unsafe_allow_html=True)